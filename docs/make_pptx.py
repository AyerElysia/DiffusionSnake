"""
生成 docs/sagittal_2d_report.pptx
风格参照 docs/report/RL_Weekly_Progress_20260709.pptx：
  - 白底
  - 内容页：深色标题栏 (1B1F2A) + 白色标题 26pt bold + 蓝色细线 (2E86DE)
  - 正文：15pt, 1B1F2A, ▸ bullet, left=0.7in top=1.6in w=11.9in h=5.4in
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H = Inches(13.333), Inches(7.5)   # 16:9

C_DARK   = RGBColor(0x1B, 0x1F, 0x2A)   # 标题栏 / 正文颜色
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE   = RGBColor(0x2E, 0x86, 0xDE)   # 细线 / 强调
C_SUB    = RGBColor(0x40, 0x40, 0x40)   # 副标题
C_DATE   = RGBColor(0x90, 0x90, 0x90)   # 日期


def _prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def _text(slide, text, x, y, w, h, size=Pt(15), bold=False,
          color=C_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
    return tb


def _content_header(slide, title):
    """深色标题栏 + 白色标题 + 蓝色细线"""
    _rect(slide, 0, 0, W, Inches(1.05), C_DARK)
    _text(slide, title,
          Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.7),
          size=Pt(26), bold=True, color=C_WHITE)
    _rect(slide, 0, Inches(1.05), W, Inches(0.04), C_BLUE)


def _body_text(slide, lines, top=Inches(1.6)):
    """正文文本框，▸ bullet，15pt，深色"""
    tb = slide.shapes.add_textbox(Inches(0.7), top, Inches(11.9), Inches(5.4))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(15)
        r.font.color.rgb = C_DARK


def _body_two_col(slide, left_title, left_lines, right_title, right_lines):
    """两列布局：左右各一个小标题 + 条目"""
    lx, rx = Inches(0.7), Inches(7.1)
    cw = Inches(5.8)

    # 左列
    _text(slide, left_title, lx, Inches(1.6), cw, Inches(0.5),
          size=Pt(16), bold=True, color=C_BLUE)
    _rect(slide, lx, Inches(2.1), cw, Inches(0.025), C_BLUE)
    body = "\n".join(left_lines)
    _text(slide, body, lx, Inches(2.2), cw, Inches(4.8),
          size=Pt(14), color=C_DARK)

    # 右列
    _text(slide, right_title, rx, Inches(1.6), cw, Inches(0.5),
          size=Pt(16), bold=True, color=C_BLUE)
    _rect(slide, rx, Inches(2.1), cw, Inches(0.025), C_BLUE)
    body2 = "\n".join(right_lines)
    _text(slide, body2, rx, Inches(2.2), cw, Inches(4.8),
          size=Pt(14), color=C_DARK)


# ─────────── Slide 1  标题页 ───────────
def slide1(prs):
    sl = _blank(prs)
    # 白底（默认），不加标题栏
    _text(sl, "矢状位椎体实例分割方案汇报",
          Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.0),
          size=Pt(32), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _text(sl, "纯2D · MoonViT特征 · Flow Matching · 时序初始化",
          Inches(1.0), Inches(3.4), Inches(11.3), Inches(0.6),
          size=Pt(20), color=C_SUB, align=PP_ALIGN.CENTER)
    _text(sl, "2026年7月",
          Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.4),
          size=Pt(14), color=C_DATE, align=PP_ALIGN.CENTER)


# ─────────── Slide 2  整体架构 ───────────
def slide2(prs):
    sl = _blank(prs)
    _content_header(sl, "整体架构")
    lines = [
        "▸  输入：单张矢状位切片（灰度图），center_repeat 复制为3通道",
        "▸  检测头：heatmap_resnet，预测26类椎体 bbox",
        "▸  特征：MoonViT 离线缓存，取 layer_18 + layer_26 各1152通道，拼接为2304-dim",
        "▸  特征对齐：LocateFeatReplacer 将2304-dim 压缩至256-dim，替换 backbone 特征",
        "▸  分割头：V4.6c DiT Flow Matching（ODE求解），Point-based 轮廓预测",
        "",
        "▸  关键原则：不引入任何跨帧信息 — 特征提取、attention、backbone 全部单帧",
        "▸  MoonViT 预训练大模型特征为只读离线缓存，不参与梯度反传",
    ]
    _body_text(sl, lines)


# ─────────── Slide 3  训练设计 ───────────
def slide3(prs):
    sl = _blank(prs)
    _content_header(sl, "训练设计")
    _body_two_col(
        sl,
        "标准2D训练（70%概率）",
        [
            "▸  输入：当前帧矢状位切片",
            "▸  初始化：由GT bbox派生的八边形初始轮廓",
            "▸  监督：当前帧GT轮廓",
            "▸  rich-state采样：在 init→GT 轨迹上随机取状态",
            "▸  模型学习：从八边形初始化精修到GT轮廓",
        ],
        "前帧轮廓初始化（30%概率）✨",
        [
            "▸  输入：仍是当前帧矢状位切片（纯2D，不变！）",
            "▸  初始化：从相邻切片mask取同类别轮廓作为初始值",
            "▸  监督：当前帧GT轮廓（与标准训练完全一致）",
            "▸  目的：模型学会「从接近GT的外部轮廓精修」",
            "▸  对齐推理：训练与推理的初始化来源保持一致",
            "▸  边界处理：无相邻帧时自动退回八边形",
        ]
    )
    _text(sl,
          "※ 纯2D：前帧信息仅用于初始轮廓，不进入特征提取，无跨帧attention，无3D结构",
          Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.5),
          size=Pt(13), color=C_BLUE)


# ─────────── Slide 4  推理流程 ───────────
def slide4(prs):
    sl = _blank(prs)
    _content_header(sl, "推理：时序轮廓传播")
    lines = [
        "▸  维护 prev_cache = {class_id: contour}，每帧推理后更新",
        "▸  case_id 发生变更时，自动清空 prev_cache（进入新病例）",
        "▸  每帧检测出 bbox 后，按 class_id 查询 prev_cache：",
        "     - 命中（该椎体上一帧有预测）→ 用前帧预测轮廓作为初始化",
        "     - 未命中（首次出现的新椎体）→ 退回 bbox 派生的八边形",
        "▸  output['py'] 为 feature 坐标系，直接缓存无需坐标变换",
        "",
        "▸  设计亮点：椎体每类唯一，按 class_id 直接匹配，无需 IoU 跟踪",
        "▸  零模型改动：复用现有 sam_i_it_py 注入路径",
        "▸  可通过 --no_temporal 关闭做消融基线对比",
    ]
    _body_text(sl, lines)


# ─────────── Slide 5  训练现状 ───────────
def slide5(prs):
    sl = _blank(prs)
    _content_header(sl, "训练现状")
    lines = [
        "▸  当前状态：正在重新提取 MoonViT 双层缓存（layer_18 + layer_26），提取完成后重启训练",
        "▸  缓存提取：约26589帧，GPU6，每分钟约95帧，预计约4.5小时完成",
        "",
        "▸  前次训练 loss 情况（约340步后停止）：",
        "     det_loss：162 → 0.88（检测头快速收敛）",
        "     diff_loss：0.5 → 0.003（flow matching loss 基本收敛）",
        "",
        "▸  已修复问题：",
        "     mask噪声碎片拆出51个instance（正常≤20）→ 每类只保留最大连通域",
        "     config key 未注册到 yacs schema → 已补充 sagittal_moonvit_fusion_mode",
        "     center_only 分支缺 patch_size/path 元数据 → 已修复",
    ]
    _body_text(sl, lines)


# ─────────── Slide 6  下一阶段思考 ───────────
def slide6(prs):
    sl = _blank(prs)
    _content_header(sl, "下一阶段思考：要不要引入3D信息？")
    _body_two_col(
        sl,
        "方向A：继续强化纯2D + 时序传播",
        [
            "▸  当前方案已有时序初始化，无需额外改动",
            "▸  训练/推理改动小，风险可控",
            "▸  等待当前训练结果，量化时序传播的收益",
            "▸  后续可调整 prev_contour_init_prob 做消融",
            "▸  代价：跨切片解剖差异大时效果受限",
        ],
        "方向B：引入真3D体素信息",
        [
            "▸  在特征层面引入相邻切片的3D上下文",
            "▸  潜在收益：对薄层/高分辨率数据更准确",
            "▸  代价：模型结构需要较大改动（大换血）",
            "▸  训练数据和显存压力显著增加",
            "▸  建议：先验证纯2D时序方案，再决定是否上3D",
        ]
    )
    _text(sl,
          "当前判断：先拿到纯2D+时序传播的定量结果，再做是否引入3D的决策",
          Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.5),
          size=Pt(13), color=C_BLUE)


# ─────────── 生成 ───────────
prs = _prs()
slide1(prs); slide2(prs); slide3(prs)
slide4(prs); slide5(prs); slide6(prs)

out = "/home/medteam/Zhrch/DiffusionSnake-12-30/docs/sagittal_2d_report.pptx"
prs.save(out)
print("saved:", out)
