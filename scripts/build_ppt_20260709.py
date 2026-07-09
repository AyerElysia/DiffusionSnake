# -*- coding: utf-8 -*-
"""Build weekly RL progress PPT."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = "/home/medteam/Zhrch/DiffusionSnake-12-30"
OUT = os.path.join(ROOT, "report", "RL_Weekly_Progress_20260709.pptx")

DARK = RGBColor(0x1B, 0x1F, 0x2A)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
ACCENT2 = RGBColor(0xE6, 0x4A, 0x4A)
GREY = RGBColor(0x5A, 0x60, 0x6E)
LIGHT = RGBColor(0xF5, 0xF6, 0xF8)
GREEN = RGBColor(0x2E, 0xA0, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def textbox(slide, l, t, w, h, text, size=18, color=DARK, bold=False,
            align=PP_ALIGN.LEFT, font="Microsoft YaHei", anchor=None, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font
    return box

def bullet_list(slide, l, t, w, h, items, size=15, color=DARK, font="Microsoft YaHei",
                 bullet_color=ACCENT, space_after=8):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            head, rest = item
            r1 = p.add_run()
            r1.text = "\u25b8 " + head
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = bullet_color
            r1.font.name = font
            if rest:
                r2 = p.add_run()
                r2.text = "  " + rest
                r2.font.size = Pt(size)
                r2.font.color.rgb = color
                r2.font.name = font
        else:
            r1 = p.add_run()
            r1.text = "\u25b8 " + item
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = font
    return box

def header_bar(slide, title, subtitle=None, page_no=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    bar.shadow.inherit = False
    textbox(slide, 0.5, 0.12, 11.5, 0.6, title, size=26, color=WHITE, bold=True)
    if subtitle:
        textbox(slide, 0.5, 0.62, 11.5, 0.35, subtitle, size=13, color=RGBColor(0xC7,0xCC,0xD6))
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05), Inches(13.333), Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    accent.shadow.inherit = False
    if page_no is not None:
        textbox(slide, 12.6, 7.12, 0.6, 0.3, str(page_no), size=11, color=GREY, align=PP_ALIGN.RIGHT)

def kpi_card(slide, l, t, w, h, value, label, value_color=ACCENT, bg=LIGHT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = RGBColor(0xE0,0xE3,0xE8)
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = value_color
    r1.font.name = "Microsoft YaHei"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(12)
    r2.font.color.rgb = GREY
    r2.font.name = "Microsoft YaHei"

# ============ Slide helpers ready. Content appended below. ============

REPORT_DIR = os.path.join(ROOT, "report")
OUT_PATH = OUT

def add_title_text(tf, text, size=32, color=DARK, bold=True):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Microsoft YaHei"

def add_body_text(tf, lines, size=16, color=DARK):
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = "Microsoft YaHei"

def add_title(slide, title):
    header_bar(slide, title)

def add_bullets(slide, items, top=1.6, font_size=15, left=0.7, width=11.9, height=5.4):
    bullet_list(slide, left, top, width, height, items, size=font_size)


# ---------------- Slide 1: 封面 ----------------
s = add_slide()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.3))
add_title_text(tb.text_frame, "强化学习精修阶段：奖励函数与动作归因问题攻关", size=32)
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.5))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "上周计划：解决强化学习的两个核心问题——奖励函数设计 与 动作归因（Credit Assignment）\n目标：解决强化学习训练中「奖励信号反馈失真」的问题"
for run in p.runs:
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
p.font.size = Pt(20)
tb3 = s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5))
add_body_text(tb3.text_frame, ["2026-07-08"], size=14, color=RGBColor(0x90,0x90,0x90))

# ---------------------------------------------------------------
# Slide 2: 背景与上周计划
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "上周计划回顾")
add_bullets(s, [
    "背景：强化学习作为精修阶段（FM + DeepSnake 结果之上再做5步几何动作精修）",
    "现象：训练信号弱、reward提升与真实IoU改善脱节，怀疑存在两个独立问题",
    "上周计划的两个攻坚方向：",
    "  ① 奖励函数设计 —— 现有 reward 是否真的贴合“精修”这个任务？",
    "  ② 动作归因（credit assignment） —— 5步动作共享一个终端reward，是否合理？",
    "目标：先把这两个问题的root cause钉死，再决定下一步扩展探索空间的方案",
], top=1.6)

# ---------------------------------------------------------------
# Slide 3: 问题1 - 动作归因诊断
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "问题①动作归因：诊断结果")
add_bullets(s, [
    "原方案：终端reward一次性算出 → 同一个标量advantage复制给5个step（等价单步RL）",
    "离线诊断（20 batch × K=8 × 5step，共4000条step-wise评分）：",
    "  · step1 的 advantage 方向与其真实贡献方向相反的比例达 25%",
    "  · step1 与 terminal 分数的相关性仅 spearman≈0.34，接近独立",
    "  · early step 真实方差只有 terminal std 的 4.2%，信号几乎被噪声淹没",
    "  · 31.25% 的探索组被 gate 完全阻断，K=8 里一条都没超过 baseline",
    "结论：不是奖励公式错，是 reward 计算层和轨迹问责层脱节，credit assignment 失效",
], top=1.6)

# ---------------------------------------------------------------
# Slide 4: 方案与对照实验
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "归因方案：full_extrap vs seq_delta")
add_bullets(s, [
    "方案设计：每一步都单独做一次“外推到完整位移”，单独打分，实现按步归因",
    "  · full_extrap（主推，GPU6）：每步外推到终点算一次reward，5步各自独立评分",
    "  · seq_delta（对照，GPU1）：用相邻步骤位移差做局部归因",
    "同一预训练起点、同样超参的对照实验结果：",
    "  · step=50起，full_extrap 就稳定领先 seq_delta，差距约0.002~0.003 IoU",
    "  · step=1000（seq_delta目标终点）：seq_delta=0.8599，extrap同期=0.8620",
    "  · seq_delta 收敛更慢更曲折（700步仅涨0.0021，extrap同等步数涨0.0016起步更快）",
    "结论：full_extrap 系统性优于 seq_delta，动作归因问题基本解决，已作为主线方案",
], top=1.6)

# ---------------------------------------------------------------
# Slide 5: FlowGRPO 路线终结
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "插曲：FlowGRPO 路线排查与终结")
add_bullets(s, [
    "并行验证了严格对齐 Flow-GRPO 论文的实现路线（内层 ODE step 的 policy 化）",
    "V1版本：old policy 快照缺失，ratio 恒为1，PPO surrogate 退化为0信号",
    "V4e版本（真实SDE latent transition logprob）：训练机制跑通，但全量验证反而低于baseline，早停",
    "V13版本（严格对齐版，多轮参数扫描，含std=0.3大噪声实验）：",
    "  100步后相对baseline仅 +0.000149 IoU，涨幅在噪声范围内，未证明有效",
    "根因：reward 来源与真正被 GRPO 更新的 policy action 没有对齐"
    "（很多正收益来自 best-of-k 搜索/蒸馏，不是被优化的动作本身）",
    "结论：FlowGRPO（inner-step policy化）路线基本放弃，不再投入",
], top=1.6)

# ---------------------------------------------------------------
# Slide 6: 问题2 - 奖励函数诊断
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "问题②奖励函数：诊断实验")
add_bullets(s, [
    "用同一checkpoint在验证集(100图/3038轮廓)上做离线诊断，分析现有4项加权reward",
    "当前权重：region 0.30 + dice 0.10 + iou 0.25 + dist 0.35（detail_score权重=0，完全闲置）",
    "核心发现：",
    "  · dice 与 iou 高度冗余（r=0.997，高质量子集r=0.9999），dice基本白费权重",
    "  · 高质量区间(iou>0.85) dice/iou 方差被压缩到几乎没有区分度",
    "  · dist_score 在高质量区间方差最大，区分度最好（是iou的约4.6倍）",
    "  · 意外发现：闲置未用的 detail_score 与真实改进量相关性 r=-0.39，强于dist(r=-0.29)",
    "  · 其中 curv_match（曲率匹配）子项相关性最强，r=-0.56",
], top=1.6)

# ---------------------------------------------------------------
# Slide 7: 奖励函数新方向 + 曲率匹配实验
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "奖励函数优化方向 + 曲率匹配验证实验")
add_bullets(s, [
    "建议新权重方向（离线诊断给出，尚未定案）：",
    "  region 0.30→0.25 ｜ dice 0.10→0（去冗余）｜ iou 保持0.25 ｜ dist 0.35→0.40 ｜ detail_score 0→0.10试验",
    "已启动验证实验：GPU4 单独只启用 curv_match 子项（其余detail子项权重=0）",
    "  与 GPU6(extrap_w1.0基线) 同一起点(step1550)对照续训",
    "当前观察（截至2250步）：两条线基本重合，curvmatch略有小幅领先趋势(0.8653 vs 0.8648)，",
    "  差距尚在噪声范围，需要更多步数才能确认曲率信号的真实增益",
], top=1.6)

# ---------------------------------------------------------------
# Slide 8: 当前训练总览曲线
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "当前训练总览：4条路线 eval_iou 曲线")
img_path = os.path.join(REPORT_DIR, "rl_curves_20260708.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.9), Inches(1.5), width=Inches(11.5))
add_bullets(s, [
    "extrap_w1.0(GPU6，主线) / curvmatch(GPU4) / delta_nsd(GPU7→GPU4) / perpoint_fmscale(GPU5)",
], top=6.9, font_size=14)

# ---------------------------------------------------------------
# Slide 9: 各路线数值小结
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "各路线关键数字小结（截至最新观测）")
rows = [
    ("路线", "当前step", "最新eval_iou", "信号状态"),
    ("extrap_w1.0 (主线,GPU6)", "2350/10000", "0.8633（峰值0.8648)", "震荡上升,是目前稳定的主线基线"),
    ("curv_match验证 (GPU4)", "2250/10000", "0.8653", "与主线基本贴合,略有领先,待更多步数确认"),
    ("delta-NSD reward (GPU7→4)", "1100/10000", "0.8659", "目前四线中数值最高,继续观察"),
    ("perpoint_fmscale (GPU5)", "900/1000", "0.8613~0.8617区间", "平台期,kl≈0,尚未验证出稳定探索空间增益"),
    ("seq_delta归因(对照,已收尾)", "1000/1000(达标)", "0.8599", "系统性弱于full_extrap 0.002~0.003,验证了主线选择正确"),
    ("flowgrpo(inner-step policy)", "多版本已终止", "—", "已判定路线失效,不再投入"),
]
top = 1.6
tbl_shape = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(top), Inches(12.2), Inches(4.6))
tbl = tbl_shape.table
widths = [Inches(3.0), Inches(1.6), Inches(2.6), Inches(5.0)]
for i, w in enumerate(widths):
    tbl.columns[i].width = w
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13) if ri > 0 else Pt(14)
            p.font.name = "Microsoft YaHei"
            if ri == 0:
                p.font.bold = True

# ---------------------------------------------------------------
# Slide 10: 可视化改进
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "配套改进：评估可视化清晰化")
add_bullets(s, [
    "问题：原默认可视化GT与匹配的Pred轮廓复用同一套调色板颜色，缩图后虚线断续看不出，无法分辨",
    "改进：新增 VIS_STYLE=clear 模式",
    "  · 线条版：GT纯绿粗实线 / Pred纯红粗虚线（加大虚线间距，缩图仍可辨）",
    "  · 填充版：GT/Pred半透明填充叠加，重叠区域自然呈现第三色，误差方向一目了然",
    "已用delta_nsd checkpoint做177样本全量评估验证：mean_iou=0.8527，与旧可视化评估结果一致，无逻辑回归",
], top=1.6)
img1 = os.path.join(REPORT_DIR, "nsd_clear_full_20260708", "montage_line_worst_median_best.png")
img2 = os.path.join(REPORT_DIR, "nsd_clear_full_20260708", "montage_fill_worst_median_best.png")
if os.path.exists(img1):
    s.shapes.add_picture(img1, Inches(0.5), Inches(4.3), height=Inches(2.9))
if os.path.exists(img2):
    s.shapes.add_picture(img2, Inches(6.9), Inches(4.3), height=Inches(2.9))

# ---------------------------------------------------------------
# Slide 11: 下一步计划
# ---------------------------------------------------------------
s = add_slide()
add_title(s, "下一步计划")
add_bullets(s, [
    "归因问题：认为已基本解决（full_extrap验证有效，是当前所有在跑路线的统一底座）",
    "奖励函数：仍在优化中",
    "  · 曲率匹配（curv_match）单项验证实验结果待定，效果不明显则考虑其他detail子项组合",
    "  · 新权重方案(降dice/升dist/试探detail)待确认后正式启用",
    "探索空间不足问题：逐点方案是目前认为最有希望的方向，但尚未突破",
    "  · 诊断出根因：逐点advantage目前是轮廓级共享标量，缺少点级定向信号",
    "  · 下一步：设计点级advantage（每个点独立的到GT距离变化量），让网络学会\"该往哪修\"",
    "运维：已重建checkpoint持久化备份守护机制，应对共享服务器外部清理风险",
], top=1.6)

prs.save(OUT_PATH)
print("SAVED:", OUT_PATH)
